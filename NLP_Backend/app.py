from flask import Flask, jsonify, request
from werkzeug.utils import secure_filename
import os
import re
import pdfplumber
from docx import Document
from pptx import Presentation
from flask_cors import CORS
from typing import List, Tuple
import requests
import numpy as np
from dotenv import load_dotenv
from collections import deque
import google.generativeai as genai
import simplejson as json

#run with python app.py not directly bcus it needs to run as a script for __name__ to be recognised
#actually run with 'python -m flask --app app run' to be extra safe...otherwise it leads to unexpected behaviour
app = Flask(__name__)
CORS(app)

load_dotenv()

app.config['UPLOAD_FOLDER'] = 'uploads/'
if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

GEMINI_API_KEY=os.getenv("GEMINI_API_KEY")
print(GEMINI_API_KEY)
# Configure Generative AI model
# genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
genai.configure(api_key=GEMINI_API_KEY)

HF_API_TOKEN = os.getenv("HUGGING_FACE_API_TOKEN")
print(HF_API_TOKEN)


rag_instance = None
conversation_history = None

def extract_text_from_ppt(ppt_file):
    text = []
    presentation = Presentation(ppt_file)
    for i, slide in enumerate(presentation.slides, 1):
        slide_text = ""
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                slide_text += shape.text + "\n"
        if slide_text.strip():
            text.append((slide_text.strip(), [f"Slide {i}"], os.path.basename(ppt_file)))
    return text

def extract_text_from_pdf(pdf_file):
    text = []
    with pdfplumber.open(pdf_file) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            page_text = page.extract_text()
            if page_text:
                text.append((page_text.strip(), [f"Page {i}"], os.path.basename(pdf_file)))
    return text

def extract_text_from_docx(docx_file):
    text = []
    doc = Document(docx_file)
    current_chunk = ""
    page_numbers = []
    page_number = 1

    for para in doc.paragraphs:
        if para._element.tag.endswith('sectPr'):
            if current_chunk.strip():
                text.append((current_chunk.strip(), page_numbers, os.path.basename(docx_file)))
            current_chunk = ""
            page_numbers = []
            page_number += 1
        else:
            if para.text.strip():
                current_chunk += para.text + "\n"
                page_numbers.append(f"Page {page_number}")
    
    if current_chunk.strip():
        text.append((current_chunk.strip(), page_numbers, os.path.basename(docx_file)))

    return text

def extract_text_from_txt(txt_file):
    with open(txt_file, 'r', encoding='utf-8') as f:
        return [(f.read(), ["Text file (no pages)"], os.path.basename(txt_file))]

def extract_text(file_path):
    if file_path.endswith('.pptx'):
        return extract_text_from_ppt(file_path)
    elif file_path.endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.endswith('.docx'):
        return extract_text_from_docx(file_path)
    elif file_path.endswith('.txt'):
        return extract_text_from_txt(file_path)
    else:
        raise ValueError("Unsupported file format")

def clean_text(text):
    cleaned_text = re.sub(r'[\x80-\xFF]+', '', text)
    cleaned_text = re.sub(r'\\[nrt]', '', cleaned_text)
    cleaned_text = re.sub(r'[^\x00-\x7F]+', ' ', cleaned_text)
    cleaned_text = re.sub(r'\s+', ' ', cleaned_text)
    return cleaned_text.strip()

def cosine_similarity(a, b):
    return np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b))

class RAG: 
    def __init__(self, text_with_refs: List[Tuple[str, List[str], str]]):
        self.chunks_with_refs = self.split_text(text_with_refs)
        self.chunks = [chunk for chunk, _, _ in self.chunks_with_refs]
        self.embeddings = self.get_embeddings(self.chunks)
        self.references = [file for _, _, file in self.chunks_with_refs]

    def split_text(self, text_with_refs: List[Tuple[str, List[str], str]], chunk_size: int = 200) -> List[Tuple[str, List[str], str]]:
        chunks_with_refs = []
        for text, ref, file in text_with_refs:
            sentences = re.split(r'(?<=[.!?])\s+', text)
            current_chunk = []
            current_length = 0
            current_pages = []
            for sentence in sentences:
                sentence = sentence.strip()
                if current_length + len(sentence.split()) > chunk_size and current_chunk:
                    chunks_with_refs.append((' '.join(current_chunk), current_pages, file))
                    current_chunk = []
                    current_length = 0
                    current_pages = []
                current_chunk.append(sentence)
                current_length += len(sentence.split())
                current_pages.extend(ref)
            if current_chunk:
                chunks_with_refs.append((' '.join(current_chunk), current_pages, file))
        return chunks_with_refs

    def get_embeddings(self, texts: List[str]) -> List[List[float]]:
        API_URL = "https://api-inference.huggingface.co/pipeline/feature-extraction/sentence-transformers/all-MiniLM-L6-v2"
        # API_URL = "https://api-inference.huggingface.co/models/sentence-transformers/all-MiniLM-L6-v2"
        headers = {"Authorization": f"Bearer {HF_API_TOKEN}"}
        
        response = requests.post(API_URL, headers=headers, json={"inputs": texts})
        response.raise_for_status()
        return response.json()

    

    def retrieve(self, query: str, k: int = 3) -> List[Tuple[str, List[str], str]]:
        query_embedding = self.get_embeddings([query])[0]
        similarities = [cosine_similarity(query_embedding, chunk_embedding) for chunk_embedding in self.embeddings]
        top_k_indices = np.argsort(similarities)[-k:][::-1]
        return [self.chunks_with_refs[i] for i in top_k_indices]

def generate_graph(response):
      
    pattern = r'\{([^}]+)\}'
    match = re.search(pattern, response.strip())

    if match:
        content = match.group(1)
        content+='}'
        content = '{' + content
        content = json.loads(content)

        x_axis = content['X axis']
        y_axis = content['Y axis']
        labels = content['label']

        return content

    return None
      

@app.route('/')
def index():
    return jsonify("Hello World")

generation_config = {
  "temperature": 1,
  "top_p": 0.95,
  "top_k": 64,
  "max_output_tokens": 8192
}

model = genai.GenerativeModel(
  model_name="gemini-1.5-flash",
  generation_config=generation_config,
  system_instruction=
            """ 
                1. You are AskYourDocs AI, an artificial intelligence assistant/chatbot designed to help users interact with and understand their uploaded documents.
                2. Your role is to provide information and assist with queries related to the content of the uploaded documents. Please focus on delivering accurate and relevant information based on the data provided.
                3. If any graphs can be generated from the document content, return the graph points.
                4. Provide graph coordinates in Python dictionary format with the keys `X axis`, `Y axis`, and `label`.
                5. The `X axis` key should contain a list of X coordinates, and the `Y axis` key should contain a list of Y coordinates.
                6. The `label` key should contain a list of labels corresponding to the X and Y coordinates.
                7. Enclose graph information within triple backticks ```Graph.
                8. Ensure all dictionary keys for graphs are enclosed in double quotes.
            """
)

PREDEFINED_PROMPT = """
You are an AI chatbot specializing in document analysis and information retrieval. Your task is to provide accurate, concise, and helpful responses based on the content extracted from various documents.

Instructions:
1. Carefully analyze the given context from the documents.
2. Provide a clear, concise, and directly relevant answer to the user's query.
3. If the query cannot be answered using the provided context, respond with: "I apologize, but I couldn't find information related to that query in the provided documents."
4. Ensure your response is well-structured, coherent, and directly addresses the user's question.
5. When appropriate, provide specific examples or elaborate to enhance understanding.
6. Avoid quoting large chunks of text verbatim. Instead, summarize and paraphrase the relevant information.
7. If the context contains conflicting information, acknowledge this and present both viewpoints with their respective sources.


Context from the documents:
{context}

User's query: {query}

Please provide your response based on these instructions:
"""

@app.route('/chat', methods=['POST'])
def chat():
    global rag_instance, conversation_history

    if 'files' in request.files:
        files = request.files.getlist('files')
        if not files or all(file.filename == '' for file in files):
            return jsonify({"error": "No files selected for uploading"}), 400

        all_text_with_refs = []
        for file in files:
            if file and file.filename:
                filename = secure_filename(file.filename)
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(file_path)

                try:
                    text_with_refs = extract_text(file_path)
                    cleaned_text_with_refs = [(clean_text(text), ref, filename) for text, ref, filename in text_with_refs]
                    all_text_with_refs.extend(cleaned_text_with_refs)
                except Exception as e:
                    return jsonify({"error": f"Error processing {filename}: {str(e)}"}), 500
                finally:
                    os.remove(file_path)

        if not all_text_with_refs:
            return jsonify({"error": "No valid content extracted from the uploaded files"}), 400

        rag_instance = RAG(all_text_with_refs)
        conversation_history = {
            'chat_session': model.start_chat(history=[]),
            'history': deque(maxlen=5)
        }

        return jsonify({"message": f"{len(files)} file(s) uploaded and processed successfully. You can now start asking questions."})

    elif 'query' in request.json:
        if not rag_instance:
            return jsonify({"error": "Please upload files before asking questions"}), 400

        query = request.json['query']
        if not query:
            return jsonify({"error": "No query provided"}), 400

        try:
            relevant_chunks = rag_instance.retrieve(query, k=5)

            if relevant_chunks:
                context = ""
                references_dict = {}
                for chunk, pages, file in relevant_chunks:
                    context += f"File: {file}, Pages: {', '.join(pages)}\n{chunk}\n\n"
                    page_numbers = set([page.split()[-1] for page in pages])
                    if file not in references_dict:
                        references_dict[file] = set()
                    references_dict[file].update(page_numbers)
                
                references = [f"{file} (Page No. {', '.join(sorted(pages, key=int))})" for file, pages in references_dict.items()]

                if not conversation_history:
                    conversation_history = {
                        'chat_session': model.start_chat(history=[]),
                        'history': deque(maxlen=5)
                    }

                chat_session = conversation_history['chat_session']

                prompt = f"""
               

                Context from the uploaded documents:
                {context}

                Query: {query}

                Please provide your response as AskYourDocs AI:
                """

                response = chat_session.send_message(prompt)

                cleaned_response = response.text.strip().replace('\n', '').replace('  ', ' ').replace('*', '').replace('\\', '')
                graph = generate_graph(cleaned_response)
                if(graph):
                    cleaned_response = ""


                conversation_history['history'].append({
                    'query': query,
                    'response': cleaned_response,
                    'references': references,
                    'graph':graph
                })

                return jsonify({
                    "response": cleaned_response,
                    "references": references,
                    'graph':graph
                })

            else:
                return jsonify({"response": "I apologize, but I couldn't find any relevant information about that in the uploaded documents. Is there something else I can help you with?", "references": []})

        except Exception as e:
            return jsonify({"error": str(e)}), 500

    else:
        return jsonify({"error": "Invalid request. Please either upload files or provide a query."}), 400
    
if __name__ == '__main__':
    app.run(debug=True)